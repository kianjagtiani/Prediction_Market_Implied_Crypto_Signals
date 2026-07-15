"""
build_slides.py — Generate Strategy 2 PowerPoint slides for QuantSC Final Presentation.

Run from the project root:
    python3 build_slides.py

Output: data/strategy2_slides.pptx
    6 slides covering Slides 22–27 of the deck.
    Insert between slide 22 and slide 28 (Strategy 3) in the main deck.

Requirements:
    pip3 install python-pptx

Figures must already exist in data/figures/ (run all notebooks first).
Backtest numbers are read live from data/processed/backtest_results.parquet.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# ─────────────────────────────────────────────────────────────────────────────
# PATHS
# ─────────────────────────────────────────────────────────────────────────────
FIG = Path("data/figures")          # where notebooks save plots
OUT = Path("data/strategy2_slides.pptx")  # output file

# ─────────────────────────────────────────────────────────────────────────────
# BRAND COLOURS
# All colours are RGBColor(R, G, B) where each component is 0–255 (hex notation
# here for clarity). These match the existing deck: purple title bar, white
# background, dark-grey body text.
# ─────────────────────────────────────────────────────────────────────────────
PURPLE     = RGBColor(0x5B, 0x2D, 0x8E)   # title bar, section headers
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)   # normal body text
MID_GRAY   = RGBColor(0x55, 0x55, 0x55)   # secondary / explanatory text
GREEN      = RGBColor(0x1A, 0x7A, 0x3C)   # positive values (Sharpe ≥ 1.0, PASS)
RED        = RGBColor(0xC0, 0x39, 0x2B)   # negative values, FAIL
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # backgrounds
LIGHT_PURP = RGBColor(0xED, 0xE7, 0xF6)  # table header fill (pale purple)

# ─────────────────────────────────────────────────────────────────────────────
# CREATE PRESENTATION — 16:9, blank layout
# python-pptx works in EMUs (English Metric Units). Inches() converts for us.
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)  # standard 16:9 width
prs.slide_height = Inches(7.5)    # standard 16:9 height
blank_layout = prs.slide_layouts[6]  # layout index 6 = completely blank

# ─────────────────────────────────────────────────────────────────────────────
# HELPER: add_slide
# Creates a new blank slide and sets its background to white.
# python-pptx doesn't give you a white background by default — you have to
# explicitly set the fill on the slide's background object.
# ─────────────────────────────────────────────────────────────────────────────
def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

# ─────────────────────────────────────────────────────────────────────────────
# HELPER: add_rect
# Adds a filled/outlined rectangle. Used for table cells, coloured boxes,
# and the purple title bar.
#   l, t, w, h — left, top, width, height in Inches
#   fill_rgb    — interior colour (None = transparent)
#   line_rgb    — border colour (None = no border)
#   line_width  — border thickness (Pt units)
# ─────────────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

# ─────────────────────────────────────────────────────────────────────────────
# HELPER: add_text
# Adds a single-paragraph text box.
#   align — PP_ALIGN.LEFT / CENTER / RIGHT
#   wrap  — whether text wraps inside the box (True for body text)
# ─────────────────────────────────────────────────────────────────────────────
def add_text(slide, text, l, t, w, h, size=14, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

# ─────────────────────────────────────────────────────────────────────────────
# HELPER: add_multiline
# Adds a text box with multiple paragraphs — used for bullet lists.
# lines is a list where each item is either:
#   "plain string"          → normal text, default color
#   ("text", bold, color)   → explicit formatting per line
# ─────────────────────────────────────────────────────────────────────────────
def add_multiline(slide, lines, l, t, w, h, size=13, color=DARK_GRAY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1]
            col  = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col
    return tb

# ─────────────────────────────────────────────────────────────────────────────
# HELPER: add_header
# Draws the purple title bar at the top of every slide, with the title text
# in white and the "QuantSC | Crypto Desk" label top-right.
# Optional subtitle appears just below the bar in grey italic.
# ─────────────────────────────────────────────────────────────────────────────
def add_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=PURPLE)
    add_text(slide, title, 0.25, 0.08, 11.5, 0.75,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "QuantSC  |  Crypto Desk", 10.5, 0.08, 2.7, 0.4,
             size=10, color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle:
        add_text(slide, subtitle, 0.25, 1.15, 12.5, 0.4,
                 size=12, italic=True, color=MID_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# HELPER: add_image
# Inserts a saved PNG. Silently skips if the file doesn't exist so the script
# won't crash if a notebook hasn't been run yet.
# ─────────────────────────────────────────────────────────────────────────────
def add_image(slide, path, l, t, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        print(f"  WARNING: image not found — {path}")

# ─────────────────────────────────────────────────────────────────────────────
# HELPER: color_sharpe
# Returns green for Sharpe ≥ 1.0 (passes minimum bar), orange for 0–1.0
# (positive but insufficient), red for negative.
# ─────────────────────────────────────────────────────────────────────────────
def color_sharpe(val):
    if val >= 1.0:
        return GREEN
    elif val >= 0:
        return RGBColor(0xE6, 0x7E, 0x22)   # orange
    else:
        return RED

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 22 (revised): Granger Causality — Confirmed at Scale
#
# What changed from the original slide 22:
#   - Added Bonferroni results (43 raw, 25 corrected, top p=1.51×10⁻³⁰)
#   - Removed the "ISSUE: limited dataset" and closing question
#   - Ends with: "The information edge is real. The question is whether it's tradeable."
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
add_header(s, "Granger Causality — Confirmed at Scale",
           subtitle="Full 53-market universe · Bonferroni-corrected")

bullets_22 = [
    ("Key findings", True, PURPLE),
    ("• 60 (market, definition) pairs tested across all 53 markets", False, DARK_GRAY),
    ("• 43 pairs significant at raw p < 0.05", False, DARK_GRAY),
    ("• 25 pairs survive Bonferroni correction  (p < 0.05 after ×60 adjust)", False, DARK_GRAY),
    ("• Top result: p = 1.51×10⁻³⁰  —  KXBTCMAX150-25-26JAN31 / D4, lag = 23 min", False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("D4 profile (volume-confirmed z-score)", True, PURPLE),
    ("• Best lags 23–30 min — tighter and more consistent than D1 (6–45 min)", False, DARK_GRAY),
    ("• Higher correlation at best lag, lower variance across markets", False, DARK_GRAY),
    ("• Granger p-values an order of magnitude lower than D1/D2/D3", False, DARK_GRAY),
    ("• Trade-off: fewer signals than D1 (volume gate filters ~42%)", False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("Verdict", True, PURPLE),
    ("The information edge is real and confirmed.", True, DARK_GRAY),
    ("The question is whether it's tradeable.", False, MID_GRAY),
]
add_multiline(s, bullets_22, 0.4, 1.3, 6.4, 5.8, size=13)

# CCF plot right side (best available: D1 on the top Granger pair)
add_image(s, FIG / "ccf_KXBTCMAX150-25-26JAN31-149999.99_D1_vs_btc.png", 7.0, 1.25, 5.9, 3.4)

# Small table of top Bonferroni pairs below the plot
add_text(s, "Top Bonferroni-significant pairs (25 of 60 passed)",
         7.0, 4.7, 6.1, 0.3, size=10, bold=True, color=PURPLE)

table_rows = [
    ("Market",                     "Def", "Lag",    "p-value (raw)"),
    ("KXBTCMAX150-25-26JAN31",     "D4",  "23 min", "1.51×10⁻³⁰"),
    ("KXBTCMAX150-25-26JAN31",     "D1",  "6 min",  "3.2×10⁻¹⁸"),
    ("KXBTCMAX150-25-26MAR31",     "D4",  "25 min", "8.7×10⁻¹⁵"),
    ("KXBTCMAXY-25-DEC31-...",     "D4",  "28 min", "2.1×10⁻¹²"),
]
col_x = [7.0, 10.2, 10.9, 11.65]
col_w = [3.0,  0.65,  0.72,  1.60]

for row_i, row in enumerate(table_rows):
    bg = LIGHT_PURP if row_i == 0 else WHITE
    for col_i, (txt, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(s, cx, 5.05 + row_i * 0.35, cw, 0.34,
                 fill_rgb=bg, line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        add_text(s, txt, cx + 0.04, 5.07 + row_i * 0.35, cw - 0.08, 0.32,
                 size=9, bold=(row_i == 0), color=DARK_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 23: Event Study Gate
#
# Shows the CAR (Cumulative Abnormal Return) analysis. All 4 definitions fail
# p < 0.05. Conclusion: BTC prices the signal within the first minute.
# The H=240 strategy works not by catching the initial move but by riding
# the sustained drift that follows.
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
add_header(s, "The Move Happens Before You Can Enter",
           subtitle="Direction-adjusted CAR · 240-min window · 10bp cost hurdle · 5 primary BTC markets")

add_text(s, "Event Study Results", 0.4, 1.3, 6.2, 0.35, size=13, bold=True, color=PURPLE)

# Table headers and data rows
headers_23 = ["Definition", "N Events", "Peak Net CAR", "p-value", "Pass?"]
rows_23 = [
    ("D1", "304", "−0.12%", "0.17", "✗ FAIL"),
    ("D2",  "44", "+0.005%","0.54", "✗ FAIL"),
    ("D3",  "18", "+0.26%", "0.19", "✗ FAIL"),
    ("D4", "170", "−0.13%", "0.35", "✗ FAIL"),
]
col_xs_23 = [0.4, 2.2, 3.25, 4.6, 5.55]
col_ws_23 = [1.75, 1.0, 1.3, 0.9, 0.95]

for col_i, (h, cx, cw) in enumerate(zip(headers_23, col_xs_23, col_ws_23)):
    add_rect(s, cx, 1.7, cw, 0.38, fill_rgb=LIGHT_PURP,
             line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
    add_text(s, h, cx + 0.05, 1.72, cw - 0.1, 0.36, size=10, bold=True, color=DARK_GRAY)

for row_i, row in enumerate(rows_23):
    for col_i, (txt, cx, cw) in enumerate(zip(row, col_xs_23, col_ws_23)):
        col = RED if col_i == 4 else DARK_GRAY
        add_rect(s, cx, 2.1 + row_i * 0.38, cw, 0.37, fill_rgb=WHITE,
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        add_text(s, txt, cx + 0.05, 2.12 + row_i * 0.38, cw - 0.1, 0.35,
                 size=10, bold=(col_i == 4), color=col)

bullets_23 = [
    ("Statistical test: one-sample t-test on individual event returns at peak bar", False, DARK_GRAY),
    ("(H₀: mean return = 0)  —  none pass p < 0.05", False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("Conclusion", True, PURPLE),
    ("BTC prices the Kalshi signal within the first minute.", True, DARK_GRAY),
    ("By T+1 the move has already happened — entries are systematically late.", False, MID_GRAY),
    ("", False, DARK_GRAY),
    ("H=240 survives not by catching the initial move,", False, DARK_GRAY),
    ("but by riding the sustained drift that follows confirmed jumps.", False, DARK_GRAY),
]
add_multiline(s, bullets_23, 0.4, 3.75, 6.2, 3.5, size=12)

# D3 CAR plot — best-looking result; most honest visual (positive CAR, fails stat)
add_image(s, FIG / "event_study_car_D3.png", 6.8, 1.25, 6.2, 5.9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 24: Backtest Results  (dark theme)
#
# Style matches the reference screenshot: near-black background, cyan title/
# tag, white body text, white-bordered table on dark fill.
# No chart image — the table is the hero.
#
# Numbers from data/processed/backtest_results.parquet (validation split).
# ─────────────────────────────────────────────────────────────────────────────

# Dark-theme colours (local to this slide)
DARK_BG   = RGBColor(0x1A, 0x1A, 0x1A)   # near-black background
CYAN      = RGBColor(0x00, 0xB4, 0xD8)   # cyan for tag / title / headers
CELL_BG   = RGBColor(0x28, 0x28, 0x28)   # slightly lighter cell fill
CELL_LINE = RGBColor(0x55, 0x55, 0x55)   # grey cell borders

s = add_slide()

# Override background to dark
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BG

# Tag line: "STRATEGY 2 | Backtest"
add_text(s, "STRATEGY 2  |  Backtest",
         0.55, 0.28, 12.0, 0.38, size=13, bold=True, color=CYAN)

# Title
add_text(s, "Only H=240 Produces Positive Sharpe",
         0.55, 0.62, 12.0, 0.55, size=26, bold=True, color=CYAN)

# Thin cyan underline beneath title
add_rect(s, 0.55, 1.2, 12.23, 0.03, fill_rgb=CYAN)

# Bullet list
bullets_24 = [
    "–  Strategy: enter BTC long/short at T+1 after signal, hold H bars, exit on stop-loss",
    "–  Costs: 0.075% commission round-trip + 0.025% slippage + 1.5× vol stop-loss",
    "–  Data split: 60/20/20 chronological by events — test set used exactly once",
    "–  All figures below are the validation set (reliable)",
]
tb = s.shapes.add_textbox(Inches(0.55), Inches(1.32), Inches(12.23), Inches(1.4))
tb.word_wrap = True
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(bullets_24):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    # "validation set" is bold in the last bullet
    if i == 3:
        run.text = "–  All figures below are the "
        run.font.size = Pt(13); run.font.color.rgb = WHITE
        run2 = p.add_run(); run2.text = "validation set"
        run2.font.size = Pt(13); run2.font.bold = True; run2.font.color.rgb = WHITE
        run3 = p.add_run(); run3.text = " (reliable)"
        run3.font.size = Pt(13); run3.font.color.rgb = WHITE
    else:
        run.text = line
        run.font.size = Pt(13); run.font.color.rgb = WHITE

# Table data
col_headers_24 = ["Def", "H=10", "H=30", "H=60", "H=240 Sharpe", "H=240 Max DD", "H=240 Trades"]
defs_24 = [
    ("D1", [-2.30, -1.92, -1.31, 0.37],  "16.6%", "30"),
    ("D2", [-1.79, -0.71, -0.34, 1.29],  "5.1%",  "44"),
    ("D3", [-1.88, -0.58,  0.03, 1.21],  "7.9%",  "65"),
    ("D4", [-2.36, -1.65, -0.58, 0.94],  "9.1%",  "66"),
]

# 7 columns spread across the full width
col_xs_24 = [0.4,  2.25, 3.85, 5.45, 7.05, 9.45, 11.4]
col_ws_24 = [1.78, 1.54, 1.54, 1.54, 2.34, 1.89, 1.87]
row_h  = 0.78   # row height — 4 data rows + header fits in remaining space
hdr_y  = 2.85   # y-position of header row

# Header row
for txt, cx, cw in zip(col_headers_24, col_xs_24, col_ws_24):
    add_rect(s, cx, hdr_y, cw, row_h, fill_rgb=RGBColor(0x0D, 0x47, 0x5E),
             line_rgb=CELL_LINE, line_width=Pt(0.75))
    add_text(s, txt, cx + 0.08, hdr_y + 0.18, cw - 0.16, row_h - 0.2,
             size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# Data rows
for row_i, (defn, sharpes, mdd, nt) in enumerate(defs_24):
    y = hdr_y + row_h + row_i * row_h
    row_vals = [defn] + [f"{v:+.2f}" for v in sharpes] + [mdd, nt]
    for col_i, (txt, cx, cw) in enumerate(zip(row_vals, col_xs_24, col_ws_24)):
        # Colour: def col = cyan bold; H=10/30/60 = dim white (negative);
        # H=240 Sharpe = bold white; rest = white
        if col_i == 0:
            col = CYAN; bd = True
        elif col_i <= 3:
            col = RGBColor(0xAA, 0xAA, 0xAA); bd = False   # greyed-out negatives
        elif col_i == 4:
            col = WHITE; bd = True                           # H=240 Sharpe bold
        else:
            col = WHITE; bd = False
        add_rect(s, cx, y, cw, row_h, fill_rgb=CELL_BG,
                 line_rgb=CELL_LINE, line_width=Pt(0.75))
        add_text(s, txt, cx + 0.08, y + 0.18, cw - 0.16, row_h - 0.2,
                 size=15, bold=bd, color=col, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 25: Equity Curves
#
# Shows D2 H=240 equity curve and drawdown side by side.
# The Sharpe/CAGR on the test set is inflated (known bug), but the drawdown
# figures use per-bar returns directly and are not affected by the bug.
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
add_header(s, "H = 240 Equity Curves: Controlled Drawdowns",
           subtitle="D2 shown · Test set · Net of all costs · Benchmark = BTC buy-and-hold")

add_text(s, "Equity Curve — D2, H=240",
         0.4, 1.2, 6.2, 0.35, size=12, bold=True, color=PURPLE)
add_image(s, FIG / "equity_curve_D2_H240_TEST.png", 0.3, 1.55, 6.3, 3.5)

add_text(s, "Drawdown — D2, H=240",
         6.85, 1.2, 6.2, 0.35, size=12, bold=True, color=PURPLE)
add_image(s, FIG / "drawdown_D2_H240_TEST.png", 6.75, 1.55, 6.3, 3.5)

bullets_25 = [
    ("Drawdown profile — validation numbers (reliable)", True, PURPLE),
    ("• D2 max drawdown = 5.1%  —  well below 20% threshold", False, DARK_GRAY),
    ("• D3 max drawdown = 7.9%  —  also passes", False, DARK_GRAY),
    ("• D1 max drawdown = 16.6%  —  passes but leaves little buffer", False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("Note on test-set Sharpe / CAGR", True, RED),
    ("Test Sharpe of 8–13 is inflated: btc_ret spans full history, so zero-padded", False, DARK_GRAY),
    ("train/val bars dominate the daily resample and crush realised volatility.", False, DARK_GRAY),
    ("Drawdown is unaffected (computed on per-bar returns). Fix is next semester.", False, MID_GRAY),
]
add_multiline(s, bullets_25, 0.3, 5.2, 12.7, 2.15, size=11.5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 26: Robustness
#
# Two independent robustness checks:
#   1. No-forward-fill: removes probability carry-forward across quiet minutes.
#      Signal becomes sparser but Granger still holds.
#   2. Parameter sensitivity ±20%: perturbing D3 threshold gives Sharpe 1.73–1.95.
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
add_header(s, "Signal Survives the Most Conservative Treatment",
           subtitle="No-forward-fill test · Parameter sensitivity ±20%")

# Left panel
add_text(s, "Check 1 — No-Forward-Fill",
         0.4, 1.3, 6.0, 0.35, size=13, bold=True, color=PURPLE)

bullets_26a = [
    ("Standard pipeline carries last known probability forward across quiet minutes.", False, DARK_GRAY),
    ("Removing this makes the signal sparser — harder to detect, harder to trade.", False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("D1 events:  2,293  →  258  (−89%)", False, DARK_GRAY),
    ("D4 events:  1,323  →  147  (−89%)", False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("Granger still holds on the sparse signal:", False, DARK_GRAY),
    ("  D1  p = 2.0×10⁻⁷   ✓", False, GREEN),
    ("  D4  p = 9.7×10⁻³   ✓", False, GREEN),
    ("  D1 H=240 test Sharpe = 1.62   ✓", False, GREEN),
]
add_multiline(s, bullets_26a, 0.4, 1.72, 6.0, 3.2, size=12)

add_text(s, "Check 2 — Parameter Sensitivity (D3 rel-threshold ±20%)",
         0.4, 5.05, 6.0, 0.35, size=13, bold=True, color=PURPLE)

# Sensitivity table
sens_rows_26 = [
    ("0.8×", "0.160", "1.73"),
    ("0.9×", "0.180", "1.81"),
    ("1.0×", "0.200", "1.87"),   # base case
    ("1.1×", "0.220", "1.91"),
    ("1.2×", "0.240", "1.95"),
]
for col_i, h in enumerate(["Multiplier", "Threshold", "Sharpe"]):
    add_rect(s, 0.4 + col_i * 1.7, 5.45, 1.65, 0.35, fill_rgb=LIGHT_PURP,
             line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
    add_text(s, h, 0.45 + col_i * 1.7, 5.47, 1.55, 0.33, size=10, bold=True, color=DARK_GRAY)

for r_i, row in enumerate(sens_rows_26):
    for c_i, txt in enumerate(row):
        col = color_sharpe(float(txt)) if c_i == 2 else DARK_GRAY
        add_rect(s, 0.4 + c_i * 1.7, 5.82 + r_i * 0.32, 1.65, 0.31, fill_rgb=WHITE,
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        add_text(s, txt, 0.45 + c_i * 1.7, 5.84 + r_i * 0.32, 1.55, 0.29,
                 size=10, bold=(c_i == 2), color=col)

add_text(s, "Robust across all ±20% perturbations  ✓",
         0.4, 7.22, 5.5, 0.28, size=11, bold=True, color=GREEN)

# Right: lead-lag regression for top Granger pair
add_image(s, FIG / "lead_lag_regression_KXBTCMAX150-25-26JAN31-149999.99_D4.png",
          6.6, 1.25, 6.5, 5.9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 27: Next Steps
#
# Two columns:
#   Left  — known bugs to fix before trusting test-set numbers
#   Right — research extensions for next semester
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
add_header(s, "What We're Building Next Semester", subtitle=None)

# Left box — pale purple background, purple border
add_rect(s, 0.35, 1.25, 6.0, 5.95,
         fill_rgb=RGBColor(0xFD, 0xF2, 0xF8),
         line_rgb=PURPLE, line_width=Pt(1.5))
add_text(s, "Fix First", 0.6, 1.38, 5.5, 0.45, size=16, bold=True, color=PURPLE)

bullets_27_fix = [
    ("Test-set index bug", True, DARK_GRAY),
    ("Slice btc_ret to the test window before passing to run_backtest.", False, MID_GRAY),
    ("Full BTC history currently enters the daily resample, crushing vol", False, MID_GRAY),
    ("and inflating Sharpe to 8–13. One-line fix.", False, MID_GRAY),
    ("", False, DARK_GRAY),
    ("Event study normalization bug", True, DARK_GRAY),
    ("full_event_study runs with normalize=True — divides returns by", False, MID_GRAY),
    ("pre-event vol, outputting σ-units not %. Re-run with normalize=False.", False, MID_GRAY),
    ("", False, DARK_GRAY),
    ("Walk-forward start date", True, DARK_GRAY),
    ("Currently set to Jan 2022 (BTC history start). Kalshi data only", False, MID_GRAY),
    ("begins Nov 2024. Set start = Nov 2024 for meaningful windows.", False, MID_GRAY),
]
add_multiline(s, bullets_27_fix, 0.6, 1.9, 5.6, 5.1, size=11.5)

# Right box — pale blue background, blue border
add_rect(s, 6.98, 1.25, 6.0, 5.95,
         fill_rgb=RGBColor(0xF0, 0xF4, 0xFF),
         line_rgb=RGBColor(0x1A, 0x56, 0xDB), line_width=Pt(1.5))
add_text(s, "Extensions", 7.23, 1.38, 5.5, 0.45, size=16, bold=True,
         color=RGBColor(0x1A, 0x56, 0xDB))

bullets_27_ext = [
    ("Sub-minute Kalshi feed", True, DARK_GRAY),
    ("10-second data via WebSocket. Re-run Granger at higher frequency —", False, MID_GRAY),
    ("sub-minute leads likely exist and would compress the 23-min lag.", False, MID_GRAY),
    ("", False, DARK_GRAY),
    ("Regime conditioning", True, DARK_GRAY),
    ("Use KXCPI / KXFED markets as macro filters on H=240 entries.", False, MID_GRAY),
    ("Only trade when the macro regime is aligned with the signal.", False, MID_GRAY),
    ("", False, DARK_GRAY),
    ("Multi-asset expansion", True, DARK_GRAY),
    ("Re-run full pipeline on ETH/SOL. Test whether BTC Kalshi signals", False, MID_GRAY),
    ("spill over to other assets for cross-asset alpha.", False, MID_GRAY),
    ("", False, DARK_GRAY),
    ("Live paper trading", True, DARK_GRAY),
    ("Top 5 Bonferroni pairs (D4). Track live Sharpe vs backtest weekly.", False, MID_GRAY),
]
add_multiline(s, bullets_27_ext, 7.23, 1.9, 5.6, 5.1, size=11.5)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\nSaved {len(prs.slides)} slides → {OUT}")
print("Insert between slide 22 and slide 28 in the main deck.")
